"""
Generates a 2-slide, 16:9 idea-submission deck for the 5G Academy 2026
(Fastweb + Vodafone track), Topic 1: AI for Network optimization.

Idea: AURA — Autonomous Root-cause & Remediation Agent (Self-Healing NOC)
Run:  .venv/Scripts/python make_slides.py  ->  AURA_Idea_Submission.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------------- palette
INK        = RGBColor(0x0A, 0x0E, 0x1A)   # near-black navy (NOC background)
PANEL      = RGBColor(0x15, 0x1B, 0x2E)   # panel navy
PANEL_2    = RGBColor(0x1E, 0x26, 0x3D)   # lighter panel
RED        = RGBColor(0xE6, 0x00, 0x00)   # Vodafone red
RED_DK     = RGBColor(0xB3, 0x00, 0x00)
CYAN       = RGBColor(0x35, 0xD0, 0xD6)   # AI accent
AMBER      = RGBColor(0xF5, 0xA6, 0x23)   # alert
GREEN      = RGBColor(0x37, 0xD6, 0x7A)   # resolved
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
MUTE       = RGBColor(0x9A, 0xA6, 0xC0)   # muted text
LINE       = RGBColor(0x2C, 0x36, 0x52)

FONT = "Segoe UI"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]


# ----------------------------------------------------------------------------- helpers
def bg(slide, color=INK):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, x, y, w, h, fill=None, line=None, line_w=1.0, shape=MSO_SHAPE.RECTANGLE,
         shadow=False, radius=None):
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    if shadow:
        el = sp._element.spPr
        ef = el.makeelement(qn('a:effectLst'), {})
        sh = ef.makeelement(qn('a:outerShdw'),
                            {'blurRad': '90000', 'dist': '40000', 'dir': '5400000', 'rotWithShape': '0'})
        clr = sh.makeelement(qn('a:srgbClr'), {'val': '000000'})
        alpha = clr.makeelement(qn('a:alpha'), {'val': '55000'})
        clr.append(alpha); sh.append(clr); ef.append(sh); el.append(ef)
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    return sp


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=2, line_spacing=1.0):
    """runs: list of paragraphs; each paragraph is list of (txt, size, color, bold, italic)."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold, italic) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.name = FONT
            r.font.color.rgb = color; r.font.bold = bold; r.font.italic = italic
    return tb


def chip(slide, x, y, w, txt, fill, fg=WHITE, size=11, h=0.32):
    rect(slide, x, y, w, h, fill=fill, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.5)
    text(slide, x, y, w, h, [[(txt, size, fg, True, False)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def line(slide, x, y, w, h, color=LINE, weight=1.5):
    cn = slide.shapes.add_connector(2, Inches(x), Inches(y), Inches(x + w), Inches(y + h))
    cn.line.color.rgb = color; cn.line.width = Pt(weight)
    cn.shadow.inherit = False
    return cn


# ============================================================================= SLIDE 1
s = prs.slides.add_slide(blank)
bg(s)

# top accent bar
rect(s, 0, 0, 13.333, 0.12, fill=RED)

# header row
chip(s, 0.6, 0.42, 3.5, "5G ACADEMY 2026  ·  TOPIC 1", PANEL_2, MUTE, 11)
chip(s, 4.25, 0.42, 3.2, "AI FOR NETWORK OPTIMIZATION", RED, WHITE, 11)
text(s, 9.4, 0.40, 3.35, 0.4, [[("FASTWEB ", 12, WHITE, True, False),
                                ("+ ", 12, MUTE, False, False),
                                ("VODAFONE", 12, RED, True, False)]],
     align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

# Title block
text(s, 0.6, 1.05, 12.1, 1.5,
     [[("AURA", 46, WHITE, True, False), ("  —  the Self-Healing NOC", 30, CYAN, True, False)],
      [("Autonomous Root-cause & Remediation Agent for Fixed–Mobile Networks", 16, MUTE, False, True)]],
     space_after=4)

# The gap sentence
text(s, 0.6, 2.55, 12.1, 0.5,
     [[("Todays networks ", 15, MUTE, False, False),
       ("detect and explain", 15, WHITE, True, False),
       (" faults. Ours ", 15, MUTE, False, False),
       ("decides and acts", 15, CYAN, True, False),
       (" — closing the loop from alarm to fix.", 15, MUTE, False, False)]])

# --- Closed-loop diagram (4 stages) ---
stages = [
    ("01", "DETECT",  "Forecast KPIs per site;\nflag anomalies early", CYAN),
    ("02", "EXPLAIN", "GenAI root-cause in\nplain language, with sources", CYAN),
    ("03", "DECIDE",  "Agent drafts the fix\n(runbook) + blast radius", AMBER),
    ("04", "ACT",     "Human approves →\nself-heal + auto report", GREEN),
]
cx, cy, cw, ch, gap = 0.6, 3.35, 2.78, 1.75, 0.32
for i, (num, name, desc, accent) in enumerate(stages):
    x = cx + i * (cw + gap)
    rect(s, x, cy, cw, ch, fill=PANEL, line=LINE, line_w=1.0,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.06, shadow=True)
    rect(s, x, cy, 0.09, ch, fill=accent)  # accent spine
    text(s, x + 0.28, cy + 0.16, cw - 0.4, 0.4,
         [[(num, 13, accent, True, False), ("   " + name, 15, WHITE, True, False)]])
    text(s, x + 0.28, cy + 0.66, cw - 0.45, 1.0,
         [[(desc, 11.5, MUTE, False, False)]], line_spacing=1.05)
    if i < 3:
        text(s, x + cw - 0.02, cy + ch/2 - 0.2, gap + 0.04, 0.4,
             [[("➜", 20, RED, True, False)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

# loop-back label
text(s, 0.6, cy + ch + 0.08, 11.5, 0.3,
     [[("closed loop  ", 11, RED, True, False),
       ("· continuous learning from every resolved incident", 11, MUTE, False, True)]])

# --- Bottom band: what it overcomes + maturity ---
by = 5.68
rect(s, 0.6, by, 7.55, 1.35, fill=PANEL, line=LINE, line_w=1.0,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
text(s, 0.85, by + 0.14, 7.1, 0.35, [[("WHAT IT OVERCOMES", 12, RED, True, False)]])
text(s, 0.85, by + 0.5, 7.15, 0.8,
     [[("• ", 12, CYAN, True, False), ("Alarm floods & slow manual root-cause  ", 11.5, WHITE, False, False),
       ("• ", 12, CYAN, True, False), ("siloed fixed vs. mobile ops", 11.5, WHITE, False, False)],
      [("• ", 12, CYAN, True, False), ("“black-box” AI engineers won’t trust  ", 11.5, WHITE, False, False),
       ("• ", 12, CYAN, True, False), ("rare faults with no training data", 11.5, WHITE, False, False)]],
     line_spacing=1.15)

# maturity callout
rect(s, 8.4, by, 4.33, 1.35, fill=RED, shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05, shadow=True)
text(s, 8.65, by + 0.16, 3.9, 0.4, [[("AUTONOMY TARGET", 11, WHITE, True, False)]])
text(s, 8.65, by + 0.5, 3.9, 0.7,
     [[("TM Forum  L2 ", 22, WHITE, True, False), ("➜", 20, WHITE, True, False),
       (" L3", 22, WHITE, True, False)]])
text(s, 8.65, by + 0.98, 3.9, 0.3, [[("from “assisted” to “conditional autonomous”", 10, WHITE, False, True)]])

# footer
line(s, 0.6, 7.18, 12.13, 0, color=LINE, weight=1)
text(s, 0.6, 7.24, 12.1, 0.25,
     [[("Team: [ your team name ]", 10, MUTE, False, False)]], align=PP_ALIGN.LEFT)
text(s, 0.6, 7.24, 12.1, 0.25,
     [[("Idea Submission · 15 July 2026", 10, MUTE, False, False)]], align=PP_ALIGN.RIGHT)


# ============================================================================= SLIDE 2
s2 = prs.slides.add_slide(blank)
bg(s2)
rect(s2, 0, 0, 13.333, 0.12, fill=RED)

# header
chip(s2, 0.6, 0.42, 2.4, "AURA · IMPACT", RED, WHITE, 11)
text(s2, 3.15, 0.40, 9.6, 0.4,
     [[("Why it matters — and why it’s feasible by October", 15, WHITE, True, False)]],
     anchor=MSO_ANCHOR.MIDDLE)

# --- three impact columns ---
cols = [
    ("SOCIAL", GREEN, "🌱", [
        "Fewer & shorter outages for millions of Italian homes and businesses",
        "Predictive cell energy-saving → lower emissions, supports net-zero goals",
        "Frees engineers from night-shift toil to higher-value work",
    ]),
    ("ECONOMIC", AMBER, "€", [
        "Cuts Mean-Time-To-Repair → direct OPEX & SLA-penalty savings",
        "Auto-written incident reports remove hours of NOC toil per event",
        "One converged platform across Fastweb + Vodafone assets",
    ]),
    ("TECHNOLOGICAL", CYAN, "◆", [
        "Moves the network up the TM Forum autonomy curve (L2→L3)",
        "GenAI closed loop with human-in-the-loop trust gating",
        "Synthetic rare-fault data → models that catch what others miss",
    ]),
]
cw2, ch2, gap2 = 3.95, 3.15, 0.24
x0 = 0.6
for i, (name, accent, icon, items) in enumerate(cols):
    x = x0 + i * (cw2 + gap2)
    rect(s2, x, 1.15, cw2, ch2, fill=PANEL, line=LINE, line_w=1.0,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05, shadow=True)
    rect(s2, x, 1.15, cw2, 0.09, fill=accent)
    # icon disk
    d = 0.62
    rect(s2, x + 0.28, 1.42, d, d, fill=PANEL_2, line=accent, line_w=1.5, shape=MSO_SHAPE.OVAL)
    text(s2, x + 0.28, 1.42, d, d, [[(icon, 20, accent, True, False)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s2, x + 1.02, 1.5, cw2 - 1.2, 0.5, [[(name, 16, WHITE, True, False)]],
         anchor=MSO_ANCHOR.MIDDLE)
    para = []
    for it in items:
        para.append([("—  ", 12, accent, True, False), (it, 11.5, MUTE, False, False)])
    text(s2, x + 0.3, 2.28, cw2 - 0.55, 1.9, para, line_spacing=1.08, space_after=7)

# --- KPI strip ---
ky = 4.65
kpis = [("↓ 40%", "target MTTR", CYAN),
        ("L2→L3", "autonomy step", RED),
        ("24/7", "self-healing NOC", GREEN),
        ("Fixed+Mobile", "one platform", AMBER)]
kw, kgap = 3.02, 0.22
for i, (big, small, accent) in enumerate(kpis):
    x = x0 + i * (kw + kgap)
    rect(s2, x, ky, kw, 0.95, fill=PANEL_2, line=LINE, line_w=1.0,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.08)
    text(s2, x, ky + 0.13, kw, 0.45, [[(big, 22, accent, True, False)]],
         align=PP_ALIGN.CENTER)
    text(s2, x, ky + 0.6, kw, 0.3, [[(small.upper(), 10, MUTE, True, False)]],
         align=PP_ALIGN.CENTER)
text(s2, x0, ky - 0.32, 12.1, 0.3,
     [[("ILLUSTRATIVE TARGETS — TO BE VALIDATED IN THE PoC", 10, MUTE, True, True)]])

# --- feasibility band ---
fy = 5.95
rect(s2, 0.6, fy, 12.13, 1.05, fill=PANEL, line=LINE, line_w=1.0,
     shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.05)
rect(s2, 0.6, fy, 0.09, 1.05, fill=GREEN)
text(s2, 0.9, fy + 0.13, 12.0, 0.35,
     [[("FEASIBLE PoC BY 15 OCTOBER  ", 12, GREEN, True, False),
       ("— we already have a working baseline", 12, MUTE, False, True)]])
text(s2, 0.9, fy + 0.5, 11.9, 0.5,
     [[("Existing forecast + anomaly pipeline (XGBoost / Isolation Forest) ", 11.5, WHITE, False, False),
       ("➜", 11.5, RED, True, False),
       (" add the GenAI Decide→Act agent + a live “talk-to-your-network” demo dashboard.", 11.5, WHITE, False, False)]],
     line_spacing=1.1)

# footer
line(s2, 0.6, 7.18, 12.13, 0, color=LINE, weight=1)
text(s2, 0.6, 7.24, 12.1, 0.25, [[("AURA · Self-Healing NOC", 10, MUTE, False, False)]])
text(s2, 0.6, 7.24, 12.1, 0.25,
     [[("Fastweb + Vodafone · 5G Academy 2026", 10, MUTE, False, False)]], align=PP_ALIGN.RIGHT)

out = "AURA_Idea_Submission.pptx"
prs.save(out)
print("Saved", out, "-", len(prs.slides._sldIdLst), "slides")
